"""Professional prezentatsiya (.pptx) yaratish moduli.

Ishlash tartibi:
  1. AI mavzu bo'yicha boy, tuzilgan kontentni JSON ko'rinishda yaratadi.
  2. python-pptx yordamida professional dizaynli .pptx fayl quriladi:
     - Bir nechta rang mavzusi (theme)
     - Turli slayd turlari: title, agenda, section, content, summary, closing
     - Footer, sahifa raqami, urg'u (accent) elementlari, "asosiy xulosa" bloki

Foydalanish:
    from bilim_ai.presentation import create_presentation
    fayl_yoli = create_presentation("Suv aylanishi", slides=10)
"""

from __future__ import annotations

import datetime
import json
import logging
import re
import tempfile
from typing import Any, Dict, List

from . import ai, config
from .prompt import presentation_prompt

logger = logging.getLogger("BilimAI.presentation")

FONT = "Calibri"


class PresentationError(Exception):
    """Prezentatsiya yaratishda yuzaga keladigan xatolar."""


# ---------------------- Rang mavzulari (themes) ----------------------
# Har biri: dark (sarlavha foni), accent (urg'u), text (qora matn),
# light (och fon), muted (kulrang).
THEMES: List[Dict[str, tuple]] = [
    {  # Ocean — chuqur ko'k
        "dark": (0x0B, 0x1F, 0x33), "accent": (0x2E, 0x8B, 0xC0),
        "text": (0x22, 0x2A, 0x33), "light": (0xEE, 0xF4, 0xF9),
        "muted": (0x8A, 0x98, 0xA6), "white": (0xFF, 0xFF, 0xFF),
    },
    {  # Emerald — zumrad yashil
        "dark": (0x0A, 0x2E, 0x24), "accent": (0x1F, 0xA9, 0x7A),
        "text": (0x20, 0x2A, 0x27), "light": (0xEC, 0xF6, 0xF2),
        "muted": (0x86, 0x97, 0x90), "white": (0xFF, 0xFF, 0xFF),
    },
    {  # Royal — siyohrang
        "dark": (0x23, 0x16, 0x40), "accent": (0x7C, 0x5C, 0xFF),
        "text": (0x26, 0x22, 0x33), "light": (0xF1, 0xEE, 0xFA),
        "muted": (0x90, 0x8A, 0xA0), "white": (0xFF, 0xFF, 0xFF),
    },
    {  # Sunset — to'q to'q sariq/qizg'ish
        "dark": (0x3A, 0x16, 0x10), "accent": (0xF0, 0x6A, 0x35),
        "text": (0x2E, 0x24, 0x22), "light": (0xFB, 0xF0, 0xEB),
        "muted": (0xA0, 0x90, 0x8A), "white": (0xFF, 0xFF, 0xFF),
    },
    {  # Corporate — kulrang/ko'k
        "dark": (0x1B, 0x26, 0x32), "accent": (0x3B, 0x82, 0xF6),
        "text": (0x1F, 0x29, 0x33), "light": (0xEF, 0xF2, 0xF6),
        "muted": (0x84, 0x90, 0x9C), "white": (0xFF, 0xFF, 0xFF),
    },
]


def _pick_theme(seed_text: str) -> Dict[str, tuple]:
    """Mavzuga qarab barqaror (lekin xilma-xil) mavzu tanlaydi."""
    s = sum(ord(c) for c in (seed_text or "BilimAI"))
    return THEMES[s % len(THEMES)]


# ---------------------- JSON tahlili ----------------------

def _extract_json(text: str) -> Dict[str, Any]:
    """AI javobidan JSON qismini ajratib oladi (markdown, ortiqcha matnga chidamli)."""
    text = (text or "").strip()
    if not text:
        raise PresentationError("AI bo'sh javob qaytardi.")

    text = re.sub(r"^```[a-zA-Z]*\s*", "", text)
    text = re.sub(r"\s*```\s*$", "", text)
    text = text.strip()

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    start = text.find("{")
    if start == -1:
        raise PresentationError("AI javobida JSON topilmadi.")
    depth = 0
    end = -1
    in_str = False
    escape = False
    for i in range(start, len(text)):
        ch = text[i]
        if escape:
            escape = False
            continue
        if ch == "\\":
            escape = True
            continue
        if ch == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                end = i + 1
                break
    if end == -1:
        raise PresentationError("AI javobida tugallanmagan JSON.")

    candidate = text[start:end]
    try:
        return json.loads(candidate)
    except json.JSONDecodeError as exc:
        cleaned = re.sub(r",(\s*[}\]])", r"\1", candidate)
        try:
            return json.loads(cleaned)
        except json.JSONDecodeError:
            raise PresentationError(
                f"AI javobidan prezentatsiya ma'lumotini o'qib bo'lmadi: {exc}"
            ) from exc


def generate_outline(topic: str, slides: int = 10) -> Dict[str, Any]:
    """AI orqali professional prezentatsiya rejasini (JSON) yaratadi (2 marta urinadi)."""
    if not config.is_configured():
        raise PresentationError(
            "AI kaliti sozlanmagan. Prezentatsiya uchun GEMINI_API_KEY yoki "
            "GROQ_API_KEY kerak."
        )

    prompt = presentation_prompt(topic, slides=slides)
    last_err: Exception | None = None
    for attempt in range(2):
        try:
            raw = ai.ask_with_system(
                prompt,
                "Sen professional prezentatsiya tuzuvchisan. Faqat va faqat to'g'ri "
                "JSON ko'rinishida javob berasan. Boshqa hech qanday matn, izoh yoki "
                "markdown qo'shma.",
            )
            data = _extract_json(raw)
            if "slides" not in data or not isinstance(data["slides"], list):
                raise PresentationError("AI javobida slaydlar topilmadi.")
            if not data["slides"]:
                raise PresentationError("AI bo'sh slaydlar ro'yxatini qaytardi.")
            return data
        except PresentationError as exc:
            last_err = exc
            logger.warning("Prezentatsiya urinishi #%d xato: %s", attempt + 1, exc)
            continue
        except Exception as exc:
            raise PresentationError(str(exc)) from exc

    raise PresentationError(
        f"Prezentatsiyani yaratib bo'lmadi (AI format xatosi). Boshqa mavzu bilan "
        f"qayta urinib ko'ring. Tafsilot: {last_err}"
    )


# ---------------------- PPTX qurish ----------------------

def build_pptx(data: Dict[str, Any], out_path: str | None = None) -> str:
    """Reja (JSON) asosida professional .pptx fayl yaratadi va yo'lini qaytaradi."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    theme = _pick_theme(str(data.get("title", "")))

    def rgb(key):
        return RGBColor(*theme[key])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def fill_bg(slide, color_key):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        rect.fill.solid()
        rect.fill.fore_color.rgb = rgb(color_key)
        rect.line.fill.background()
        rect.shadow.inherit = False
        slide.shapes._spTree.remove(rect._element)
        slide.shapes._spTree.insert(2, rect._element)
        return rect

    def add_rect(slide, left, top, width, height, color_key):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(color_key)
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def textbox(slide, left, top, width, height, anchor=None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        return tf

    def set_para(p, text, size, color_key, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, space_after=8):
        p.text = text
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = rgb(color_key)
        return p

    def footer(slide, page_no, total, title_txt):
        # Pastki ajratuvchi chiziq
        add_rect(slide, Inches(0.6), Inches(6.95), Inches(12.13), Pt(1.2), "light")
        tf = textbox(slide, Inches(0.6), Inches(7.0), Inches(9), Inches(0.4))
        set_para(tf.paragraphs[0], f"BilimAI  •  {title_txt}"[:70], 10, "muted")
        tf2 = textbox(slide, Inches(11.0), Inches(7.0), Inches(1.73), Inches(0.4))
        set_para(tf2.paragraphs[0], f"{page_no} / {total}", 10, "muted",
                 align=PP_ALIGN.RIGHT)

    title_text = str(data.get("title", "Prezentatsiya")).strip() or "Prezentatsiya"
    subtitle_text = str(data.get("subtitle", "")).strip()
    author = str(data.get("author", "BilimAI")).strip() or "BilimAI"
    today = datetime.date.today().strftime("%d.%m.%Y")

    raw_slides = [s for s in data.get("slides", []) if isinstance(s, dict)]
    # title slaydni AI ro'yxatidan chiqarib tashlaymiz (biz o'zimiz quramiz)
    raw_slides = [s for s in raw_slides if str(s.get("type", "")).lower() != "title"]
    total_pages = len(raw_slides) + 1  # + title slayd

    # ---------- 1) TITLE SLAYD ----------
    slide = prs.slides.add_slide(blank)
    fill_bg(slide, "dark")
    # Urg'u bloki (chap tomonda vertikal)
    add_rect(slide, 0, Inches(2.4), Inches(0.22), Inches(2.6), "accent")
    tf = textbox(slide, Inches(0.9), Inches(2.3), Inches(11.4), Inches(2.8),
                 anchor=MSO_ANCHOR.MIDDLE)
    set_para(tf.paragraphs[0], title_text, 48, "white", bold=True, space_after=10)
    if subtitle_text:
        p = tf.add_paragraph()
        set_para(p, subtitle_text, 22, "accent", space_after=4)
    # Pastki ma'lumot
    tf2 = textbox(slide, Inches(0.95), Inches(6.5), Inches(11), Inches(0.6))
    set_para(tf2.paragraphs[0], f"{author}  •  {today}", 14, "muted")

    # ---------- Qolgan slaydlar ----------
    for idx, item in enumerate(raw_slides, start=2):
        stype = str(item.get("type", "content")).lower()
        heading = str(item.get("heading", "")).strip()
        bullets = [str(b).strip() for b in (item.get("bullets") or []) if str(b).strip()]
        takeaway = str(item.get("takeaway", "")).strip()

        slide = prs.slides.add_slide(blank)

        # --- SECTION (bo'lim ajratuvchi) ---
        if stype == "section":
            fill_bg(slide, "dark")
            add_rect(slide, Inches(0.9), Inches(3.05), Inches(1.6), Inches(0.14),
                     "accent")
            tf = textbox(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(2.0),
                         anchor=MSO_ANCHOR.TOP)
            set_para(tf.paragraphs[0], heading or "Bo'lim", 40, "white", bold=True)
            # katta tartib raqami (fon)
            tfn = textbox(slide, Inches(0.85), Inches(1.4), Inches(4), Inches(1.4))
            set_para(tfn.paragraphs[0], f"{idx - 1:02d}", 60, "accent", bold=True)
            continue

        # --- CLOSING (yakuniy) ---
        if stype == "closing":
            fill_bg(slide, "dark")
            add_rect(slide, 0, Inches(2.6), Inches(0.22), Inches(2.2), "accent")
            tf = textbox(slide, Inches(0.9), Inches(2.6), Inches(11.4), Inches(2.2),
                         anchor=MSO_ANCHOR.MIDDLE)
            set_para(tf.paragraphs[0], heading or "Rahmat!", 46, "white", bold=True)
            sub = bullets[0] if bullets else "Savollaringiz bormi?"
            p = tf.add_paragraph()
            set_para(p, sub, 22, "accent")
            tf2 = textbox(slide, Inches(0.95), Inches(6.5), Inches(11), Inches(0.6))
            set_para(tf2.paragraphs[0], f"{author}  •  {today}", 14, "muted")
            continue

        # --- CONTENT / AGENDA / SUMMARY (och fonli) ---
        fill_bg(slide, "white")
        # chap urg'u chizig'i
        add_rect(slide, 0, 0, Inches(0.22), SH, "accent")
        # sarlavha
        head_tf = textbox(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(1.0))
        set_para(head_tf.paragraphs[0], heading or "Slayd", 30, "dark", bold=True)
        # sarlavha tagidagi urg'u chizig'i
        add_rect(slide, Inches(0.72), Inches(1.45), Inches(2.0), Pt(3.5), "accent")

        is_agenda = stype == "agenda"
        is_summary = stype == "summary"

        body_top = Inches(1.9)
        body = textbox(slide, Inches(0.8), body_top, Inches(11.8), Inches(4.6))
        first = True
        for i, b in enumerate(bullets, start=1):
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            # marker run
            if is_agenda:
                marker = f"{i}.  "
            elif is_summary:
                marker = "✓  "
            else:
                marker = "▸  "
            run_m = p.add_run()
            run_m.text = marker
            run_m.font.name = FONT
            run_m.font.size = Pt(20)
            run_m.font.bold = True
            run_m.font.color.rgb = rgb("accent")
            run_t = p.add_run()
            run_t.text = b
            run_t.font.name = FONT
            run_t.font.size = Pt(20)
            run_t.font.color.rgb = rgb("text")
            p.space_after = Pt(16)

        # takeaway (asosiy xulosa) bloki
        if takeaway:
            box = add_rect(slide, Inches(0.8), Inches(6.05), Inches(11.7),
                           Inches(0.75), "light")
            box.line.fill.background()
            ttf = box.text_frame
            ttf.word_wrap = True
            ttf.margin_left = Inches(0.2)
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = ttf.paragraphs[0]
            run_m = pp.add_run()
            run_m.text = "💡  "
            run_m.font.size = Pt(15)
            run_t = pp.add_run()
            run_t.text = takeaway
            run_t.font.name = FONT
            run_t.font.size = Pt(15)
            run_t.font.bold = True
            run_t.font.italic = True
            run_t.font.color.rgb = rgb("dark")

        footer(slide, idx, total_pages, title_text)

    if out_path is None:
        tmp = tempfile.NamedTemporaryFile(
            delete=False, suffix=".pptx", prefix="bilimai_"
        )
        out_path = tmp.name
        tmp.close()

    prs.save(out_path)
    return out_path


def create_presentation(topic: str, slides: int = 10, out_path: str | None = None) -> str:
    """Mavzu bo'yicha to'liq professional prezentatsiya yaratadi. Fayl yo'lini qaytaradi."""
    topic = (topic or "").strip()
    if not topic:
        raise PresentationError("Prezentatsiya uchun mavzu bo'sh bo'lishi mumkin emas.")
    slides = max(6, min(slides, 16))
    outline = generate_outline(topic, slides=slides)
    return build_pptx(outline, out_path=out_path)
